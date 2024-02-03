from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

# Create a new Presentation object
prs = Presentation()

# Define the title and content for the first two slides
slides_content = [
    ("Data Ingestion and Transformation", "An overview of Domain 1, focusing on the critical aspects of data ingestion and transformation within AWS services. This includes understanding throughput, latency, patterns, and the distinction between streaming and batch data ingestion methods."),
    ("Perform Data Ingestion", "Key Knowledge:\n- Throughput and latency of AWS data ingestion services\n- Data ingestion patterns\n- Streaming vs. batch data ingestion methods\n- Replayability and state management in ingestion pipelines\n\nSkills:\n- Utilizing streaming sources: Kinesis, MSK, etc.\n- Batch sources handling: S3, Glue, EMR\n- Ingestion configuration and security considerations.")
]

# Use a title and content layout for the slides
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Choosing a title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation
presentation_path = "/mnt/data/data_ingestion_presentation.pptx"
prs.save(presentation_path)

presentation_path
